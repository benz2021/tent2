import streamlit as st
from PIL import Image, ImageDraw, ImageFont
import os
import io
import base64
import requests
import tempfile
from pathlib import Path

class ImageLayoutApp:
    def __init__(self):
        # ตัวแปรสำหรับเก็บสถานะ
        self.image_paths = []
        self.preview_images = []
        self.image_files = []
        
        # ค่าเริ่มต้น (หน่วยเซนติเมตร)
        self.img_width = 8.0
        self.img_height = 8.0
        self.spacing = 1.0
        self.margin = 2.0
        self.paper_orientation = "แนวตั้ง"
        self.rotate_top = True
        self.auto_fit = False
        
        # DPI สำหรับแปลง cm -> pixel
        self.DPI = 300
        
    def cm_to_px(self, cm):
        return int(cm * self.DPI / 2.54)
    
    def get_github_files(self, repo_url, path=""):
        """ดึงรายการไฟล์จาก GitHub repository"""
        try:
            # แปลง GitHub URL เป็น API URL
            if "github.com" in repo_url:
                # แปลง URL เช่น https://github.com/username/repo -> https://api.github.com/repos/username/repo
                api_url = repo_url.replace("github.com", "api.github.com/repos")
                if not api_url.endswith("/contents"):
                    api_url = api_url.rstrip("/") + "/contents"
                
                if path:
                    api_url = api_url + "/" + path
            else:
                # ถ้าเป็น raw URL หรืออื่นๆ
                return []
            
            response = requests.get(api_url)
            if response.status_code == 200:
                files = response.json()
                image_files = []
                for file in files:
                    if file['type'] == 'file':
                        ext = os.path.splitext(file['name'])[1].lower()
                        if ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']:
                            image_files.append({
                                'name': file['name'],
                                'download_url': file['download_url'],
                                'path': file['path']
                            })
                    elif file['type'] == 'dir':
                        # ดึงไฟล์จาก subdirectory
                        sub_files = self.get_github_files(repo_url, file['path'])
                        image_files.extend(sub_files)
                
                return image_files
            else:
                st.error(f"ไม่สามารถเชื่อมต่อ GitHub: {response.status_code}")
                return []
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาด: {str(e)}")
            return []
    
    def load_image_from_url(self, url):
        """โหลดภาพจาก URL"""
        try:
            response = requests.get(url)
            if response.status_code == 200:
                img = Image.open(io.BytesIO(response.content))
                return img
            else:
                st.error(f"ไม่สามารถโหลดภาพ: {url}")
                return None
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาด: {str(e)}")
            return None
    
    def generate_preview(self):
        """สร้างภาพตัวอย่าง"""
        if not self.image_files:
            st.warning("กรุณาเลือกไฟล์ภาพจาก GitHub")
            return
        
        try:
            self.preview_images.clear()
            
            # ดึงค่าจาก session state
            img_width_cm = st.session_state.get('img_width', self.img_width)
            img_height_cm = st.session_state.get('img_height', self.img_height)
            spacing_cm = st.session_state.get('spacing', self.spacing)
            margin_cm = st.session_state.get('margin', self.margin)
            paper_orientation = st.session_state.get('paper_orientation', self.paper_orientation)
            rotate_top = st.session_state.get('rotate_top', self.rotate_top)
            auto_fit = st.session_state.get('auto_fit', self.auto_fit)
            
            # ขนาดกระดาษ A4 (cm)
            if paper_orientation == "แนวตั้ง":
                paper_w_cm, paper_h_cm = 21.0, 29.7
            else:
                paper_w_cm, paper_h_cm = 29.7, 21.0
            
            # แปลงเป็น pixel
            paper_w = self.cm_to_px(paper_w_cm)
            paper_h = self.cm_to_px(paper_h_cm)
            margin = self.cm_to_px(margin_cm)
            
            progress_bar = st.progress(0)
            
            # อ่านและสร้างภาพแต่ละแผ่น
            for idx, file_info in enumerate(self.image_files):
                # โหลดภาพจาก URL
                img = self.load_image_from_url(file_info['download_url'])
                if img is None:
                    continue
                
                # ถ้าเลือกปรับขนาดอัตโนมัติ
                if auto_fit:
                    # คำนวณขนาดให้พอดี (2 ภาพซ้อนกัน)
                    max_h = (paper_h - 2 * margin - self.cm_to_px(spacing_cm)) / 2
                    max_w = paper_w - 2 * margin
                    
                    # คำนวณอัตราส่วน
                    img_ratio = img.width / img.height
                    
                    # ปรับขนาดตามความสูงที่จำกัด
                    img_h = int(max_h)
                    img_w = int(max_h * img_ratio)
                    
                    # ถ้ากว้างเกินไป ให้ปรับตามความกว้าง
                    if img_w > max_w:
                        img_w = int(max_w)
                        img_h = int(max_w / img_ratio)
                else:
                    # ใช้ขนาดที่ผู้ใช้กำหนด
                    img_w = self.cm_to_px(img_width_cm)
                    img_h = self.cm_to_px(img_height_cm)
                
                # ตรวจสอบว่าภาพใหญ่เกินกระดาษหรือไม่
                total_height = 2 * img_h + self.cm_to_px(spacing_cm)
                if total_height > paper_h - 2 * margin:
                    st.warning(f"ภาพที่ {idx+1} ({file_info['name']}) ใหญ่เกินไป กรุณาลดขนาดหรือเปิด Auto Fit")
                    continue
                
                # สร้างกระดาษเปล่า
                paper = Image.new('RGB', (paper_w, paper_h), 'white')
                
                # ปรับขนาดภาพต้นฉบับ
                img_resized = img.resize((img_w, img_h), Image.Resampling.LANCZOS)
                
                # ---- ภาพบน (หมุน 180 องศา) ----
                top_img = img_resized.copy()
                if rotate_top:
                    top_img = top_img.rotate(180)
                
                # ---- ภาพล่าง (ไม่หมุน) ----
                bottom_img = img_resized.copy()
                
                # ---- คำนวณตำแหน่งกึ่งกลาง ----
                total_height = 2 * img_h + self.cm_to_px(spacing_cm)
                start_y = (paper_h - total_height) // 2
                start_x = (paper_w - img_w) // 2
                
                # ---- วางภาพ ----
                # ภาพบน
                paper.paste(top_img, (start_x, start_y))
                
                # ภาพล่าง
                bottom_y = start_y + img_h + self.cm_to_px(spacing_cm)
                paper.paste(bottom_img, (start_x, bottom_y))
                
                # เก็บภาพ
                self.preview_images.append(paper)
                
                # อัพเดท progress
                progress_bar.progress((idx + 1) / len(self.image_files))
            
            progress_bar.empty()
            
            if self.preview_images:
                st.success(f"✅ สร้างตัวอย่างสำเร็จ {len(self.preview_images)} แผ่น")
            else:
                st.warning("ไม่สามารถสร้างภาพได้ กรุณาตรวจสอบไฟล์ที่เลือก")
                
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาด: {str(e)}")

    def create_pptx(self):
        """สร้างไฟล์ PowerPoint (.pptx) โดยแต่ละภาพตัวอย่าง = 1 สไลด์
        ขนาดสไลด์จะตรงกับขนาดกระดาษ A4 ที่เลือก (แนวตั้ง/แนวนอน)
        คืนค่าเป็น BytesIO หรือ None หากเกิดข้อผิดพลาด
        """
        if not self.preview_images:
            return None

        try:
            from pptx import Presentation
            from pptx.util import Cm
        except ImportError:
            st.warning("กรุณาติดตั้ง python-pptx: pip install python-pptx")
            return None

        paper_orientation = st.session_state.get('paper_orientation', self.paper_orientation)
        if paper_orientation == "แนวตั้ง":
            paper_w_cm, paper_h_cm = 21.0, 29.7
        else:
            paper_w_cm, paper_h_cm = 29.7, 21.0

        try:
            prs = Presentation()
            # กำหนดขนาดสไลด์ให้เท่ากับขนาดกระดาษ A4
            prs.slide_width = Cm(paper_w_cm)
            prs.slide_height = Cm(paper_h_cm)

            blank_layout = prs.slide_layouts[6]  # เลย์เอาต์เปล่า (ไม่มี placeholder)

            for img in self.preview_images:
                slide = prs.slides.add_slide(blank_layout)

                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                img_bytes.seek(0)

                # วางภาพให้เต็มสไลด์พอดีกับขนาดกระดาษ
                slide.shapes.add_picture(
                    img_bytes,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )

            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes

        except Exception as e:
            st.error(f"เกิดข้อผิดพลาดขณะสร้าง PowerPoint: {str(e)}")
            return None

    def create_download_buttons(self):
        """สร้างปุ่มดาวน์โหลด"""
        if not self.preview_images:
            return
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # ดาวน์โหลด PNG แต่ละแผ่น
            for idx, img in enumerate(self.preview_images):
                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                img_bytes.seek(0)
                
                st.download_button(
                    label=f"📥 ดาวน์โหลด แผ่นที่ {idx+1} (PNG)",
                    data=img_bytes,
                    file_name=f"page_{idx+1:03d}.png",
                    mime="image/png"
                )
        
        with col2:
            # ดาวน์โหลด PDF รวมทุกแผ่น
            try:
                import img2pdf
                
                # แปลงทุกภาพเป็น bytes
                image_bytes_list = []
                for img in self.preview_images:
                    img_bytes = io.BytesIO()
                    img.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    image_bytes_list.append(img_bytes.getvalue())
                
                # สร้าง PDF
                pdf_bytes = img2pdf.convert(image_bytes_list)
                
                st.download_button(
                    label="📄 ดาวน์โหลด PDF (หลายหน้า)",
                    data=pdf_bytes,
                    file_name="output.pdf",
                    mime="application/pdf"
                )
            except ImportError:
                st.warning("กรุณาติดตั้ง img2pdf: pip install img2pdf")

        with col3:
            # ดาวน์โหลด PowerPoint รวมทุกแผ่น (1 แผ่น = 1 สไลด์)
            pptx_bytes = self.create_pptx()
            if pptx_bytes is not None:
                st.download_button(
                    label="📊 ดาวน์โหลด PowerPoint (.pptx)",
                    data=pptx_bytes,
                    file_name="output.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

def main():
    st.set_page_config(
        page_title="จัดเรียงภาพบนกระดาษ A4",
        page_icon="🖼️",
        layout="wide"
    )
    
    st.title("📐 จัดเรียงภาพบนกระดาษ A4 (1 แผ่น = 2 ภาพ)")
    st.markdown("---")
    
    # สร้าง instance ของแอพ
    if 'app' not in st.session_state:
        st.session_state.app = ImageLayoutApp()
    
    app = st.session_state.app
    
    # Sidebar สำหรับการตั้งค่า
    with st.sidebar:
        st.header("⚙️ การตั้งค่า")
        
        # ส่วนของ GitHub
        st.subheader("📂 ดึงไฟล์จาก GitHub")
        
        github_url = st.text_input(
            "URL GitHub Repository:",
            placeholder="https://github.com/username/repository",
            key="github_url"
        )
        
        if st.button("🔍 ดึงภาพจาก GitHub"):
            if github_url:
                with st.spinner("กำลังดึงไฟล์จาก GitHub..."):
                    files = app.get_github_files(github_url)
                    if files:
                        app.image_files = files
                        app.image_paths = [f['name'] for f in files]
                        st.success(f"พบภาพ {len(files)} ไฟล์")
                    else:
                        st.warning("ไม่พบไฟล์ภาพใน GitHub repository นี้")
            else:
                st.warning("กรุณาใส่ URL ของ GitHub repository")
        
        # แสดงรายการไฟล์ที่พบ
        if app.image_files:
            st.subheader(f"📋 ไฟล์ที่พบ ({len(app.image_files)} ไฟล์)")
            for file in app.image_files:
                st.text(f"• {file['name']}")
        
        st.markdown("---")
        
        # การตั้งค่าขนาด
        st.subheader("📏 ขนาดภาพ (เซนติเมตร)")
        
        st.session_state.img_width = st.number_input(
            "ความกว้างภาพ:",
            min_value=0.5,
            max_value=20.0,
            value=app.img_width,
            step=0.5,
            key="img_width_input"
        )
        
        st.session_state.img_height = st.number_input(
            "ความสูงภาพ:",
            min_value=0.5,
            max_value=20.0,
            value=app.img_height,
            step=0.5,
            key="img_height_input"
        )
        
        st.session_state.spacing = st.number_input(
            "ระยะห่างระหว่างภาพบน-ล่าง:",
            min_value=0.0,
            max_value=5.0,
            value=app.spacing,
            step=0.5,
            key="spacing_input"
        )
        
        st.session_state.margin = st.number_input(
            "ขอบกระดาษ:",
            min_value=0.0,
            max_value=5.0,
            value=app.margin,
            step=0.5,
            key="margin_input"
        )
        
        st.markdown("---")
        
        # ตัวเลือกเพิ่มเติม
        st.subheader("🎯 ตัวเลือกเพิ่มเติม")
        
        st.session_state.paper_orientation = st.selectbox(
            "การวางแนวกระดาษ:",
            ["แนวตั้ง", "แนวนอน"],
            index=0,
            key="orientation_select"
        )
        
        st.session_state.rotate_top = st.checkbox(
            "หมุนภาพบน 180 องศา",
            value=app.rotate_top,
            key="rotate_checkbox"
        )
        
        st.session_state.auto_fit = st.checkbox(
            "ปรับขนาดภาพอัตโนมัติให้พอดีกระดาษ",
            value=app.auto_fit,
            key="autofit_checkbox"
        )
        
        st.markdown("---")
        
        # ปุ่มสร้างตัวอย่าง
        if st.button("🔍 สร้างตัวอย่าง", type="primary", use_container_width=True):
            app.generate_preview()
    
    # ส่วนแสดงผลหลัก
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("🖼️ ตัวอย่างภาพ")
        
        if app.preview_images:
            # แสดงภาพตัวอย่าง
            preview_idx = st.session_state.get('preview_idx', 0)
            
            if preview_idx >= len(app.preview_images):
                preview_idx = len(app.preview_images) - 1
                st.session_state.preview_idx = preview_idx
            
            # แสดงภาพ
            st.image(app.preview_images[preview_idx], caption=f"แผ่นที่ {preview_idx+1}/{len(app.preview_images)}")
            
            # ปุ่มนำทาง
            if len(app.preview_images) > 1:
                col_prev, col_page, col_next = st.columns([1, 2, 1])
                with col_prev:
                    if st.button("◀ ก่อนหน้า") and preview_idx > 0:
                        st.session_state.preview_idx = preview_idx - 1
                        st.rerun()
                with col_next:
                    if st.button("ถัดไป ▶") and preview_idx < len(app.preview_images) - 1:
                        st.session_state.preview_idx = preview_idx + 1
                        st.rerun()
        else:
            st.info("💡 ยังไม่มีภาพตัวอย่าง กรุณาคลิก 'สร้างตัวอย่าง'")
    
    with col2:
        st.header("💾 ดาวน์โหลด")
        app.create_download_buttons()
    
    # แสดงข้อมูลสถานะ
    if app.image_files:
        st.markdown("---")
        col_info1, col_info2, col_info3 = st.columns(3)
        with col_info1:
            st.metric("จำนวนไฟล์ภาพ", len(app.image_files))
        with col_info2:
            st.metric("จำนวนแผ่นที่สร้าง", len(app.preview_images))
        with col_info3:
            orientation = st.session_state.get('paper_orientation', 'แนวตั้ง')
            st.metric("การวางแนว", orientation)

if __name__ == "__main__":
    main()
